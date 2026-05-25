"""Generate a .pptx file from a Markdown slide outline.

Markdown format expected:

    # Slide Title
    - Bullet point 1
    - Bullet point 2
      - Sub-bullet
    > Speaker notes here

Each H1 starts a new slide. Bullets become content. Blockquotes become speaker notes.

Requires ``python-pptx`` (optional dependency). Falls back gracefully if not installed.
"""
from __future__ import annotations

import re
from pathlib import Path


class SlidesError(RuntimeError):
    def __init__(self, code: str, message: str):
        super().__init__(message)
        self.code = code


def _parse_outline(md: str) -> list[dict]:
    """Parse markdown into list of slide dicts: {title, bullets, notes}."""
    slides: list[dict] = []
    current: dict | None = None

    for line in md.splitlines():
        if line.startswith("# "):
            if current:
                slides.append(current)
            current = {"title": line[2:].strip(), "bullets": [], "notes": ""}
        elif current is None:
            continue
        elif line.startswith("> "):
            current["notes"] += line[2:].strip() + "\n"
        elif re.match(r"^\s*[-*]\s", line):
            bullet = re.sub(r"^\s*[-*]\s", "", line)
            indent = len(line) - len(line.lstrip())
            current["bullets"].append({"text": bullet, "level": min(indent // 2, 2)})
        elif line.strip():
            current["bullets"].append({"text": line.strip(), "level": 0})

    if current:
        slides.append(current)
    return slides


def generate_pptx(outline_md: str, output_path: Path) -> Path:
    """Generate a .pptx file from markdown outline. Returns the output path."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise SlidesError(
            "PPTX_NOT_AVAILABLE",
            "python-pptx is not installed. Install with: pip install python-pptx"
        )

    slides_data = _parse_outline(outline_md)
    if not slides_data:
        raise SlidesError("EMPTY_OUTLINE", "No slides found in the outline markdown.")

    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        slide.shapes.title.text = slide_data["title"]

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()

        for i, bullet in enumerate(slide_data["bullets"]):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = bullet["text"]
            p.level = bullet["level"]
            p.font.size = Pt(18 - bullet["level"] * 2)

        if slide_data["notes"].strip():
            slide.notes_slide.notes_text_frame.text = slide_data["notes"].strip()

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    return output_path


def outline_to_markdown_slides(outline_md: str) -> str:
    """Convert outline markdown to a structured slide-by-slide summary (no pptx needed)."""
    slides = _parse_outline(outline_md)
    parts = []
    for i, s in enumerate(slides, 1):
        bullets = "\n".join(f"{'  ' * b['level']}- {b['text']}" for b in s["bullets"])
        part = f"### 第{i}页: {s['title']}\n{bullets}"
        if s["notes"].strip():
            part += f"\n> 备注: {s['notes'].strip()}"
        parts.append(part)
    return "\n\n".join(parts)
